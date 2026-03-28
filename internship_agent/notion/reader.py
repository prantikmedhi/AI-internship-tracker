"""
Notion content readers - read pages, databases, and file attachments.
"""
import io
import os
import requests
from typing import Optional

from internship_agent.notion.resolver import resolve_page_id


def read_page_blocks(page_id: str, client) -> str:
    """
    Reads text content from a standard Notion page (not database rows).

    Returns all text blocks with their IDs and types.

    Args:
        page_id: Page UUID or name
        client: Notion client

    Returns:
        Formatted text content with block metadata
    """
    if not client:
        return "Notion client not initialized."

    page_id = resolve_page_id(page_id, client)

    try:
        results = client.blocks.children.list(block_id=page_id)
        output = []

        for block in results.get("results", []):
            block_type = block["type"]
            block_id = block["id"]

            # Extract text from readable block types
            if block_type in [
                "paragraph", "heading_1", "heading_2", "heading_3",
                "bulleted_list_item", "numbered_list_item",
                "to_do", "toggle", "quote"
            ]:
                try:
                    text_arr = block[block_type].get("rich_text", [])
                    text = "".join([t.get("plain_text", "") for t in text_arr])
                    output.append(f"[{block_id} | {block_type}] {text}")
                except (KeyError, TypeError):
                    pass

        if not output:
            return (
                "No readable text blocks found on this page. "
                "(If this is a Database, use read_database instead)."
            )

        return "\n".join(output)

    except Exception as e:
        return f"Error reading page {page_id}: {str(e)}"


def read_database(database_id: str, client) -> str:
    """
    Reads all rows and columns from a Notion database.

    Shows column headers and includes blank cells marked as "[blank]".

    Args:
        database_id: Database UUID or name
        client: Notion client

    Returns:
        Formatted table with all rows and columns
    """
    if not client:
        return "Notion client not initialized."

    database_id = resolve_page_id(database_id, client)

    try:
        # Get schema to show headers
        db_schema = client.databases.retrieve(database_id=database_id)
        db_props = db_schema.get("properties", {})
        column_names = list(db_props.keys())

        # Query database rows
        results = client.databases.query(database_id=database_id)
        pages = results.get("results", [])

        if not pages:
            return "This database is empty."

        output = []

        # Display column headers
        header_line = " | ".join(column_names)
        output.append(f"COLUMNS: {header_line}")
        output.append("─" * len(header_line))

        # Display each row
        for i, page in enumerate(pages[:15]):
            props = page.get("properties", {})
            row_data = []

            # Iterate through columns in order
            for col_name in column_names:
                prop_val = props.get(col_name, {})
                val_str = ""
                p_type = prop_val.get("type", "")

                # Extract value based on property type
                if p_type == "title" and prop_val.get("title"):
                    val_str = prop_val["title"][0]["plain_text"]
                elif p_type == "rich_text" and prop_val.get("rich_text"):
                    val_str = "".join([
                        t["plain_text"] for t in prop_val["rich_text"]
                    ])
                elif p_type == "number" and prop_val.get("number") is not None:
                    val_str = str(prop_val["number"])
                elif p_type == "select" and prop_val.get("select"):
                    val_str = prop_val["select"]["name"]
                elif p_type == "multi_select" and prop_val.get("multi_select"):
                    val_str = ", ".join([
                        s["name"] for s in prop_val["multi_select"]
                    ])
                elif p_type == "date" and prop_val.get("date"):
                    val_str = prop_val["date"].get("start", "")
                elif p_type == "checkbox":
                    val_str = "Yes" if prop_val.get("checkbox") else "No"
                elif p_type == "url" and prop_val.get("url"):
                    val_str = prop_val["url"]

                # Show blank cells as [blank]
                row_data.append(val_str if val_str else "[blank]")

            row_str = " | ".join(row_data)
            output.append(f"Row {i+1} [ID: {page['id']}]: {row_str}")

        return "\n".join(output)

    except Exception as e:
        return f"Error querying database {database_id}: {str(e)}"


def read_file_attachments(page_id: str, client) -> str:
    """
    Downloads and reads all file attachments from a Notion page.

    Supported formats: PDF, DOCX, XLSX, CSV, PPTX

    Args:
        page_id: Page UUID or name
        client: Notion client

    Returns:
        Extracted text content from all files
    """
    if not client:
        return "Notion client not initialized."

    page_id = resolve_page_id(page_id, client)

    try:
        blocks = client.blocks.children.list(block_id=page_id)
        output = []

        for block in blocks.get("results", []):
            block_type = block.get("type", "")

            if block_type not in ("file", "pdf", "image"):
                continue

            file_data = block.get(block_type, {})
            url = (
                file_data.get("file", {}).get("url") or
                file_data.get("external", {}).get("url")
            )
            name = file_data.get("name", "unknown_file")

            if not url:
                continue

            # Download file
            print(f"Downloading: {name} from {url[:80]}...")
            resp = requests.get(url, timeout=30)
            resp.raise_for_status()

            ext = os.path.splitext(name.split("?")[0])[1].lower()
            content_bytes = resp.content
            text = _extract_file_content(ext, content_bytes, name)

            output.append(f"=== File: {name} ===\n{text[:8000]}")

        if not output:
            return "No file attachments found on this page."

        return "\n\n".join(output)

    except Exception as e:
        return f"Error reading files from page {page_id}: {str(e)}"


def _extract_file_content(ext: str, content_bytes: bytes, filename: str) -> str:
    """Extract text from various file formats."""
    try:
        if ext == ".pdf":
            import fitz
            pdf = fitz.open(stream=content_bytes, filetype="pdf")
            return "\n".join(page.get_text() for page in pdf)

        elif ext == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext in (".xlsx", ".xls"):
            import pandas as pd
            df_dict = pd.read_excel(io.BytesIO(content_bytes), sheet_name=None)
            text = ""
            for sheet, df in df_dict.items():
                text += f"\n--- Sheet: {sheet} ---\n{df.to_string()}\n"
            return text

        elif ext == ".csv":
            import pandas as pd
            df = pd.read_csv(io.BytesIO(content_bytes))
            return df.to_string()

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"\n--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            # Try UTF-8 decode for txt or unknown
            return content_bytes.decode("utf-8", errors="ignore")

    except Exception as e:
        return f"Could not parse {filename}: {e}"
